import io
import os
import re

from pptx import Presentation
from pptx.util import Pt

try:
    from image_scrapper import downloader
except ImportError:
    from .image_scrapper import downloader


async def generate_ppt_prompt(language, emotion_type, slide_length, topic):
    message = f"""Create an {language} language outline for a {emotion_type} slideshow presentation on the topic of {topic} which is {slide_length} slides long. 
Make sure it is {slide_length} slides long.

You are allowed to use the following slide types:

Slide types:
Title Slide - (Title, Subtitle)
Content Slide - (Title, Content)
Image Slide - (Title, Content, Image)
Thanks Slide - (Title)

Put this tag before the Title Slide: [L_TS]
Put this tag before the Content Slide: [L_CS]
Put` this tag before the Image Slide: [L_IS]
Put this tag before the Thanks Slide: [L_THS]

Put this tag after each Slide: [SLIDEBREAK]

For example:
[L_TS]
[TITLE]Mount Everest: The Highest Peak in the World[/TITLE]

[SLIDEBREAK]

[L_IS]
[TITLE]Facts about Mount Everest[/TITLE]
[CONTENT]• It is 8,848 meters (29,029 ft) high above sea level
• First successfully climbed by Sir Edmund Hillary and Tenzing Norgay on May 29, 1953
• Over 300 climbers have died attempting to scale the mountain[/CONTENT]
[IMAGE]Mount Everest[/IMAGE]

[SLIDEBREAK]

Put this tag before the Title: [TITLE]
Put this tag after the Title: [/TITLE]
Put this tag before the Subtitle: [SUBTITLE]
Put this tag after the Subtitle: [/SUBTITLE]
Put this tag before the Content: [CONTENT]
Put this tag after the Content: [/CONTENT]
Put this tag before the Image: [IMAGE]
Put this tag after the Image: [/IMAGE]

Elaborate on the Content, provide as much information as possible.
You put a [/CONTENT] at the end of the Content.
Pay attention to the language of presentation - {language}.
Each image should be described in general by a set of keywords, such as "Mount Everest Sunset" or "Niagara Falls Rainbow".
Do not write Image in Content tag.
Do not reply as if you are talking about the slideshow itself. (ex. "Include pictures here about...")
Do not include any special characters (?, !, ., :, ) in the Title.
Do not include any additional information in your response and stick to the format."""

    return message


async def generate_ppt(answer, template,time):
    template = os.path.join("src", "ai_generator", "presentation_templates", f"{template}.pptx")
    print(template,"template")
    root = Presentation(template)

    # """ Ref for slide types:
    # 0 -> title and subtitle
    # 1 -> title and content
    # 2 -> section header
    # 3 -> two content
    # 4 -> Comparison
    # 5 -> Title only
    # 6 -> Blank
    # 7 -> Content with caption
    # 8 -> Pic with caption
    # """

    async def delete_all_slides():
        for i in range(len(root.slides) - 1, -1, -1):
            r_id = root.slides._sldIdLst[i].rId
            root.part.drop_rel(r_id)
            del root.slides._sldIdLst[i]

    async def create_title_slide(title, subtitle):
        layout = root.slide_layouts[0]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(30)

    async def create_section_header_slide(title):
        layout = root.slide_layouts[2]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(25)
    def set_font_size_for_all_paragraphs(text_frame, font_size_pt):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size_pt)
    async def create_title_and_content_slide(title, content):
        layout = root.slide_layouts[1]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(25)
        slide.placeholders[1].text = content
        content_shape = slide.placeholders[1]
        content_text_frame = content_shape.text_frame
        set_font_size_for_all_paragraphs(content_text_frame, 15)
       

    async def create_title_and_content_and_image_slide(title, content, image_query):
        layout = root.slide_layouts[8]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(22)
        slide.placeholders[2].text = content
        content_shape = slide.placeholders[2]
        content_text_frame = content_shape.text_frame
        content_text_frame.paragraphs[0].font.size = Pt(15)
        content_text_frame.paragraphs[1].font.size = Pt(15)
        # content_text_frame.paragraphs[2].font.size = Pt(20)
        # content_text_frame.paragraphs[3].font.size = Pt(20)
        try:
            image_data = await downloader.download(image_query, limit=1, adult_filter_off=True, timeout=15,
                                                   filter="+filterui:aspect-wide+filterui:imagesize-wallpaper+filterui:photo-photo")
            slide.placeholders[1].insert_picture(io.BytesIO(image_data))
        except Exception:
            pass

    async def find_text_in_between_tags(text, start_tag, end_tag):
        start_pos = text.find(start_tag)
        end_pos = text.find(end_tag)
        result = []
        while start_pos > -1 and end_pos > -1:
            text_between_tags = text[start_pos + len(start_tag):end_pos]
            result.append(text_between_tags)
            start_pos = text.find(start_tag, end_pos + len(end_tag))
            end_pos = text.find(end_tag, start_pos)
        res1 = "".join(result)
        res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
        if len(result) > 0:
            return res2
        else:
            return ""

    async def search_for_slide_type(text):
        tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
        found_text = next((s for s in tags if s in text), None)
        return found_text

    async def parse_response(reply):
        list_of_slides = reply.split("[SLIDEBREAK]")
        for slide in list_of_slides:
            slide_type = await search_for_slide_type(slide)
            match slide_type:
                case ("[L_TS]"):
                    await create_title_slide(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                             await find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
                case ("[L_CS]"):
                    await create_title_and_content_slide("".join(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                                                         "".join(await find_text_in_between_tags(str(slide), "[CONTENT]", "[/CONTENT]")))
                case ("[L_IS]"):
                    await create_title_and_content_and_image_slide("".join(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                                                                   "".join(await find_text_in_between_tags(str(slide), "[CONTENT]", "[/CONTENT]")),
                                                                   "".join(await find_text_in_between_tags(str(slide), "[IMAGE]", "[/IMAGE]")))
                case ("[L_THS]"):
                    await create_section_header_slide("".join(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))

    async def find_title():
        return root.slides[0].shapes.title.text

    await delete_all_slides()
    await parse_response(answer)
    buffer = io.BytesIO()
    root.save(buffer)
    pptx_bytes = buffer.getvalue()
    pptx_title = f"{await find_title()}.pptx"
    # print(f"done {pptx_title}")

    with open(f"output2{time}.pptx", "wb") as f:
        f.write(pptx_bytes)
    print(f"Presentation saved as output.pptx")
    return pptx_bytes, pptx_title


async def generate_ppt(answer, template, time):
    template = os.path.join("src", "ai_generator", "presentation_templates", f"{template}.pptx")
    print(template, "template")
    root = Presentation(template)

    async def delete_all_slides():
        for i in range(len(root.slides) - 1, -1, -1):
            r_id = root.slides._sldIdLst[i].rId
            root.part.drop_rel(r_id)
            del root.slides._sldIdLst[i]

    async def create_title_slide(title, author):
        layout = root.slide_layouts[0]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"Author: {author}"
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(30)

    async def create_plans_slide(plans):
        layout = root.slide_layouts[1]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = "Reja:"
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(25)
        limited_plans = plans[:3]
        plans_content = "\n".join([f"{plan}" for i, plan in enumerate(limited_plans )])
        slide.placeholders[1].text = plans_content
        content_shape = slide.placeholders[1]
        content_text_frame = content_shape.text_frame
        set_font_size_for_all_paragraphs(content_text_frame, 15)

    async def create_title_slide(title, subtitle):
        layout = root.slide_layouts[0]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(30)

    async def create_section_header_slide(title):
        layout = root.slide_layouts[2]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(25)
    def set_font_size_for_all_paragraphs(text_frame, font_size_pt):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size_pt)
    async def create_title_and_content_slide(title, content):
        layout = root.slide_layouts[1]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(25)
        slide.placeholders[1].text = content
        content_shape = slide.placeholders[1]
        content_text_frame = content_shape.text_frame
        set_font_size_for_all_paragraphs(content_text_frame, 15)
       

    async def create_title_and_content_and_image_slide(title, content, image_query):
        layout = root.slide_layouts[8]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(22)
        slide.placeholders[2].text = content
        content_shape = slide.placeholders[2]
        content_text_frame = content_shape.text_frame
        content_text_frame.paragraphs[0].font.size = Pt(15)
        content_text_frame.paragraphs[1].font.size = Pt(15)
        # content_text_frame.paragraphs[2].font.size = Pt(20)
        # content_text_frame.paragraphs[3].font.size = Pt(20)
        try:
            image_data = await downloader.download(image_query, limit=1, adult_filter_off=True, timeout=15,
                                                   filter="+filterui:aspect-wide+filterui:imagesize-wallpaper+filterui:photo-photo")
            slide.placeholders[1].insert_picture(io.BytesIO(image_data))
        except Exception:
            pass

    def set_font_size_for_all_paragraphs(text_frame, font_size_pt):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size_pt)

    async def find_text_in_between_tags(text, start_tag, end_tag):
        start_pos = text.find(start_tag)
        end_pos = text.find(end_tag)
        result = []
        while start_pos > -1 and end_pos > -1:
            text_between_tags = text[start_pos + len(start_tag):end_pos]
            result.append(text_between_tags)
            start_pos = text.find(start_tag, end_pos + len(end_tag))
            end_pos = text.find(end_tag, start_pos)
        res1 = "".join(result)
        res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
        if len(result) > 0:
            return res2
        else:
            return ""

    async def search_for_slide_type(text):
        tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]", "[L_PN]"]
        found_text = next((s for s in tags if s in text), None)
        return found_text

    async def parse_response(reply):
        list_of_slides = reply.split("[SLIDEBREAK]")
        plans = []
        author = ""
        title_slide_created = False
        for slide in list_of_slides:
            slide_type = await search_for_slide_type(slide)
            match slide_type:
                case "[L_TS]":
                    title = await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")
                    author = await find_text_in_between_tags(str(slide), "[AUTHOR]", "[/AUTHOR]")
                    await create_title_slide(title, author)
                    title_slide_created = True
                case "[L_PN]":
                    plan = await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")
                    plans.append(plan)

        if  plans:
            # Insert the plans slide at index 1 (right after the title slide)
         await create_plans_slide(plans)
         
        for slide in list_of_slides:
            slide_type = await search_for_slide_type(slide)
            match slide_type:
                # case "[L_TS]":
                #     title = await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")
                #     author = await find_text_in_between_tags(str(slide), "[AUTHOR]", "[/AUTHOR]")
                #     await create_title_slide(title, author)
                #     title_slide_created = True
                # case "[L_PN]":
                #     plan = await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")
                #     plans.append(plan)
                    
                    
                case ("[L_CS]"):
                    await create_title_and_content_slide("".join(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                                                         "".join(await find_text_in_between_tags(str(slide), "[CONTENT]", "[/CONTENT]")))
                case ("[L_IS]"):
                    await create_title_and_content_and_image_slide("".join(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                                                                   "".join(await find_text_in_between_tags(str(slide), "[CONTENT]", "[/CONTENT]")),
                                                                   "".join(await find_text_in_between_tags(str(slide), "[IMAGE]", "[/IMAGE]")))
                case ("[L_THS]"):
                    await create_section_header_slide("".join(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))
        
        
        
        # After parsing all slides, create a single slide for plans
        

    async def find_title():
        return root.slides[0].shapes.title.text

    await delete_all_slides()
    await parse_response(answer)
    buffer = io.BytesIO()
    root.save(buffer)
    pptx_bytes = buffer.getvalue()
    pptx_title = f"{await find_title()}.pptx"

    with open(f"output2{time}.pptx", "wb") as f:
        f.write(pptx_bytes)
    print(f"Presentation saved as output.pptx")
    return pptx_bytes, pptx_title

import sys
import asyncio

async def test():
    text = sys.stdin.read()
    templates = (text.split("{{")[1]).strip().split("_")[0]
    timestamp = (text.split("{{")[1]).strip().split("_")[1]
    txt =(text.split("{{")[0]).strip()
    
    output, error = await generate_ppt(txt, templates,timestamp)
    print(output, error,"Salom")

if __name__ == "__main__":
    asyncio.run(test())